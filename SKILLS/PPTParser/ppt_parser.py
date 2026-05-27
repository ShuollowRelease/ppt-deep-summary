"""
PPTParser - PPT 深度解析组件

按照 skill.md 规范实现 PPTX 文件的结构化解析与数据清洗，
输出稳定、标准化的 JSON 和 Markdown 中间数据。
"""

import json
from pathlib import Path
from typing import Any
from dataclasses import dataclass, field, asdict

from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.autoshape import Shape
from pptx.util import Inches


@dataclass
class TextItem:
    level: int = 0
    text: str = ""


@dataclass
class TableData:
    rows: list[list[str]] = field(default_factory=list)
    row_count: int = 0
    col_count: int = 0


@dataclass
class ImageData:
    filename: str = ""
    content_type: str = ""
    width: float = 0.0
    height: float = 0.0


@dataclass
class HyperlinkData:
    text: str = ""
    url: str = ""


@dataclass
class ChartData:
    chart_type: str = ""
    title: str = ""


@dataclass
class SlideData:
    slide_number: int = 0
    title: str | None = None
    layout_name: str = ""
    texts: list[TextItem] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)
    tables: list[TableData] = field(default_factory=list)
    images: list[ImageData] = field(default_factory=list)
    hyperlinks: list[HyperlinkData] = field(default_factory=list)
    charts: list[ChartData] = field(default_factory=list)
    raw_markdown: str = ""


@dataclass
class ParsedPresentation:
    presentation_title: str = ""
    slide_count: int = 0
    slides: list[SlideData] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


class PPTParserError(Exception):
    pass


def validate_file(pptx_path: Path) -> None:
    if not pptx_path.exists():
        raise PPTParserError(f"文件不存在: {pptx_path}")
    if pptx_path.suffix.lower() not in ['.pptx', '.ppt']:
        raise PPTParserError(f"不支持的文件格式: {pptx_path.suffix}，仅支持 .pptx")


def extract_texts_from_frame(text_frame) -> list[TextItem]:
    items: list[TextItem] = []
    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            items.append(TextItem(level=paragraph.level or 0, text=text))
    return items


def extract_title_from_shape(shape: Shape) -> str | None:
    if not shape.has_text_frame:
        return None
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            return text
    return None


def extract_texts_recursively(shape) -> list[TextItem]:
    results: list[TextItem] = []
    if isinstance(shape, GroupShape):
        for child in shape.shapes:
            results.extend(extract_texts_recursively(child))
    elif shape.has_text_frame:
        results.extend(extract_texts_from_frame(shape.text_frame))
    return results


def extract_table(table) -> TableData:
    data = TableData()
    data.row_count = len(table.rows)
    data.col_count = len(table.columns)
    for row in table.rows:
        row_data = [cell.text.strip() for cell in row.cells]
        data.rows.append(row_data)
    return data


def extract_image(shape) -> ImageData | None:
    if not hasattr(shape, 'image'):
        return None
    try:
        image = shape.image
        return ImageData(
            filename=image.filename,
            content_type=image.content_type,
            width=shape.width / 914400 if shape.width else 0,
            height=shape.height / 914400 if shape.height else 0
        )
    except Exception:
        return None


def extract_hyperlink(shape) -> HyperlinkData | None:
    if not shape.has_text_frame:
        return None
    links: list[HyperlinkData] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.hyperlink and run.hyperlink.address:
                links.append(HyperlinkData(
                    text=run.text.strip(),
                    url=run.hyperlink.address
                ))
    return links[0] if links else None


def extract_hyperlinks_from_shape(shape) -> list[HyperlinkData]:
    links: list[HyperlinkData] = []
    if isinstance(shape, GroupShape):
        for child in shape.shapes:
            links.extend(extract_hyperlinks_from_shape(child))
    elif shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    if run.hyperlink and run.hyperlink.address:
                        links.append(HyperlinkData(
                            text=run.text.strip(),
                            url=run.hyperlink.address
                        ))
                except AttributeError:
                    pass
    return links


def extract_chart(shape) -> ChartData | None:
    if not shape.has_chart:
        return None
    try:
        chart = shape.chart
        chart_type = str(chart.chart_type) if chart.chart_type else "unknown"
        title = chart.chart_title.text_frame.text if chart.has_title else ""
        return ChartData(chart_type=chart_type, title=title)
    except Exception:
        return None


def extract_notes(slide) -> list[str]:
    notes_text: list[str] = []
    if not slide.has_notes_slide:
        return notes_text
    notes_slide = slide.notes_slide
    if notes_slide and notes_slide.notes_text_frame:
        for paragraph in notes_slide.notes_text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                notes_text.append(text)
    return notes_text


def process_slide_shapes(slide) -> tuple[str | None, list[TextItem], list[TableData], list[ImageData], list[HyperlinkData], list[ChartData]]:
    title: str | None = None
    texts: list[TextItem] = []
    tables: list[TableData] = []
    images: list[ImageData] = []
    hyperlinks: list[HyperlinkData] = []
    charts: list[ChartData] = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                extracted_title = extract_title_from_shape(shape)
                if extracted_title:
                    title = extracted_title
            else:
                texts.extend(extract_texts_recursively(shape))

        if shape.has_table:
            tables.append(extract_table(shape.table))

        image = extract_image(shape)
        if image:
            images.append(image)

        hyperlinks.extend(extract_hyperlinks_from_shape(shape))

        chart = extract_chart(shape)
        if chart:
            charts.append(chart)

    return title, texts, tables, images, hyperlinks, charts


def generate_markdown(slide_data: SlideData) -> str:
    lines: list[str] = []

    lines.append(f"# Slide {slide_data.slide_number}")
    lines.append("")

    if slide_data.layout_name:
        lines.append(f"**Layout:** {slide_data.layout_name}")
        lines.append("")

    if slide_data.title:
        lines.append("## " + slide_data.title)
        lines.append("")

    if slide_data.texts:
        for item in slide_data.texts:
            indent = "  " * item.level
            prefix = "- " if item.level > 0 else ""
            lines.append(f"{indent}{prefix}{item.text}")
        lines.append("")

    if slide_data.notes:
        lines.append("### Notes")
        for note in slide_data.notes:
            lines.append(f"> {note}")
        lines.append("")

    if slide_data.tables:
        for idx, table in enumerate(slide_data.tables):
            lines.append(f"### Table {idx + 1}")
            if table.rows:
                header = " | ".join(table.rows[0])
                lines.append(f"| {header} |")
                lines.append("|" + " --- |" * len(table.rows[0]))
                for row in table.rows[1:]:
                    row_str = " | ".join(row)
                    lines.append(f"| {row_str} |")
            lines.append("")

    if slide_data.images:
        lines.append("### Images")
        for img in slide_data.images:
            lines.append(f"- {img.filename} ({img.width:.1f}x{img.height:.1f} inches)")
        lines.append("")

    if slide_data.hyperlinks:
        lines.append("### Hyperlinks")
        for link in slide_data.hyperlinks:
            lines.append(f"- [{link.text}]({link.url})")
        lines.append("")

    if slide_data.charts:
        lines.append("### Charts")
        for chart in slide_data.charts:
            lines.append(f"- Type: {chart.chart_type}")
            if chart.title:
                lines.append(f"  Title: {chart.title}")
        lines.append("")

    return "\n".join(lines)


def parse_pptx(pptx_path: str | Path) -> ParsedPresentation:
    pptx_path = Path(pptx_path)
    validate_file(pptx_path)

    result = ParsedPresentation()
    warnings: list[str] = []

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        raise PPTParserError(f"无法打开 PPT 文件: {e}")

    result.presentation_title = pptx_path.stem
    result.slide_count = len(prs.slides)

    for idx, slide in enumerate(prs.slides, start=1):
        try:
            slide_data = SlideData(slide_number=idx)

            layout = slide.slide_layout
            if layout and layout.name:
                slide_data.layout_name = layout.name

            title, texts, tables, images, hyperlinks, charts = process_slide_shapes(slide)

            slide_data.title = title
            slide_data.texts = texts
            slide_data.tables = tables
            slide_data.images = images
            slide_data.hyperlinks = hyperlinks
            slide_data.charts = charts

            try:
                slide_data.notes = extract_notes(slide)
            except Exception as e:
                warnings.append(f"Slide {idx} notes extraction failed: {e}")

            slide_data.raw_markdown = generate_markdown(slide_data)
            result.slides.append(slide_data)

        except Exception as e:
            warnings.append(f"Slide {idx} parsing failed: {e}")
            result.slides.append(SlideData(
                slide_number=idx,
                raw_markdown=f"# Slide {idx}\n\n[Parsing Failed: {e}]"
            ))

    result.warnings = warnings
    return result


def to_dict(parsed: ParsedPresentation) -> dict[str, Any]:
    return {
        "presentation_title": parsed.presentation_title,
        "slide_count": parsed.slide_count,
        "slides": [
            {
                "slide_number": s.slide_number,
                "title": s.title,
                "layout_name": s.layout_name,
                "texts": [{"level": t.level, "text": t.text} for t in s.texts],
                "notes": s.notes,
                "tables": [{"rows": t.rows, "row_count": t.row_count, "col_count": t.col_count} for t in s.tables],
                "images": [{"filename": i.filename, "content_type": i.content_type, "width": i.width, "height": i.height} for i in s.images],
                "hyperlinks": [{"text": h.text, "url": h.url} for h in s.hyperlinks],
                "charts": [{"chart_type": c.chart_type, "title": c.title} for c in s.charts],
                "raw_markdown": s.raw_markdown
            }
            for s in parsed.slides
        ],
        "warnings": parsed.warnings if parsed.warnings else []
    }


def to_json(parsed: ParsedPresentation, indent: int = 2) -> str:
    return json.dumps(to_dict(parsed), ensure_ascii=False, indent=indent)


def to_markdown(parsed: ParsedPresentation) -> str:
    lines: list[str] = []
    lines.append(f"# {parsed.presentation_title}")
    lines.append("")
    lines.append(f"Total Slides: {parsed.slide_count}")
    lines.append("")
    lines.append("---")
    lines.append("")

    for slide in parsed.slides:
        lines.append(slide.raw_markdown)
        lines.append("---")
        lines.append("")

    if parsed.warnings:
        lines.append("## Warnings")
        for w in parsed.warnings:
            lines.append(f"- {w}")

    return "\n".join(lines)


def parse(pptx_path: str | Path, json_output: str | Path | None = None, md_output: str | Path | None = None) -> ParsedPresentation:
    parsed = parse_pptx(pptx_path)

    if json_output:
        json_path = Path(json_output)
        json_path.parent.mkdir(parents=True, exist_ok=True)
        json_path.write_text(to_json(parsed), encoding='utf-8')

    if md_output:
        md_path = Path(md_output)
        md_path.parent.mkdir(parents=True, exist_ok=True)
        md_path.write_text(to_markdown(parsed), encoding='utf-8')

    return parsed


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法: python ppt_parser.py <pptx文件路径> [输出JSON路径] [输出Markdown路径]")
        sys.exit(1)

    pptx_file = sys.argv[1]
    json_file = sys.argv[2] if len(sys.argv) > 2 else None
    md_file = sys.argv[3] if len(sys.argv) > 3 else None

    try:
        result = parse(pptx_file, json_file, md_file)
        print(to_json(result))
        if result.warnings:
            print("\nWarnings:", file=sys.stderr)
            for w in result.warnings:
                print(f"  - {w}", file=sys.stderr)
    except PPTParserError as e:
        print(json.dumps({"status": "error", "reason": str(e)}))
        sys.exit(1)
