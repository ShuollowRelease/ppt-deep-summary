import json
from pathlib import Path
from typing import Any
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.autoshape import Shape


@dataclass
class TextItem:
    level: int = 0
    text: str = ""


@dataclass
class TableData:
    rows: list[list[str]] = field(default_factory=list)
    row_count: int = 0
    col_count: int = 0


@dataclass
class ImageData:
    filename: str = ""
    content_type: str = ""
    width: float = 0.0
    height: float = 0.0


@dataclass
class HyperlinkData:
    text: str = ""
    url: str = ""


@dataclass
class ChartData:
    chart_type: str = ""
    title: str = ""


@dataclass
class SlideData:
    slide_number: int = 0
    title: str | None = None
    layout_name: str = ""
    texts: list[TextItem] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)
    tables: list[TableData] = field(default_factory=list)
    images: list[ImageData] = field(default_factory=list)
    hyperlinks: list[HyperlinkData] = field(default_factory=list)
    charts: list[ChartData] = field(default_factory=list)
    raw_markdown: str = ""


@dataclass
class ParsedPresentation:
    presentation_title: str = ""
    slide_count: int = 0
    slides: list[SlideData] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


class PPTParserError(Exception):
    pass


def validate_file(pptx_path: Path) -> None:
    if not pptx_path.exists():
        raise PPTParserError(f"文件不存在: {pptx_path}")
    if pptx_path.suffix.lower() not in ['.pptx', '.ppt']:
        raise PPTParserError(f"不支持的文件格式: {pptx_path.suffix}")


def extract_texts_from_frame(text_frame) -> list[TextItem]:
    return [TextItem(level=p.level or 0, text=p.text.strip()) for p in text_frame.paragraphs if p.text.strip()]


def extract_texts_recursively(shape) -> list[TextItem]:
    if isinstance(shape, GroupShape):
        return [t for child in shape.shapes for t in extract_texts_recursively(child)]
    return extract_texts_from_frame(shape.text_frame) if shape.has_text_frame else []


def extract_table(table) -> TableData:
    return TableData(
        rows=[[cell.text.strip() for cell in row.cells] for row in table.rows],
        row_count=len(table.rows),
        col_count=len(table.columns)
    )


def extract_image(shape) -> ImageData | None:
    if not hasattr(shape, 'image'):
        return None
    try:
        img = shape.image
        return ImageData(filename=img.filename, content_type=img.content_type,
                         width=shape.width / 914400 if shape.width else 0,
                         height=shape.height / 914400 if shape.height else 0)
    except Exception:
        return None


def extract_hyperlinks(shape) -> list[HyperlinkData]:
    links: list[HyperlinkData] = []
    if isinstance(shape, GroupShape):
        return [h for child in shape.shapes for h in extract_hyperlinks(child)]
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                try:
                    if run.hyperlink and run.hyperlink.address:
                        links.append(HyperlinkData(text=run.text.strip(), url=run.hyperlink.address))
                except AttributeError:
                    pass
    return links


def extract_chart(shape) -> ChartData | None:
    if not shape.has_chart:
        return None
    try:
        chart = shape.chart
        return ChartData(chart_type=str(chart.chart_type) if chart.chart_type else "unknown",
                         title=chart.chart_title.text_frame.text if chart.has_title else "")
    except Exception:
        return None


def extract_notes(slide) -> list[str]:
    if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame:
        return []
    return [p.text.strip() for p in slide.notes_slide.notes_text_frame.paragraphs if p.text.strip()]


def process_slide(slide) -> tuple[str | None, list[TextItem], list[TableData], list[ImageData], list[HyperlinkData], list[ChartData]]:
    title = None
    texts, tables, images, hyperlinks, charts = [], [], [], [], []

    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        title = p.text.strip()
                        break
            else:
                texts.extend(extract_texts_recursively(shape))
        if shape.has_table:
            tables.append(extract_table(shape.table))
        img = extract_image(shape)
        if img:
            images.append(img)
        hyperlinks.extend(extract_hyperlinks(shape))
        chart = extract_chart(shape)
        if chart:
            charts.append(chart)

    return title, texts, tables, images, hyperlinks, charts


def generate_markdown(s: SlideData) -> str:
    lines = [f"# Slide {s.slide_number}", ""]
    if s.layout_name:
        lines.extend([f"**Layout:** {s.layout_name}", ""])
    if s.title:
        lines.extend(["## " + s.title, ""])
    if s.texts:
        lines.extend([f"{'  ' * t.level}{'- ' if t.level > 0 else ''}{t.text}" for t in s.texts] + [""])
    if s.notes:
        lines.extend(["### Notes"] + [f"> {n}" for n in s.notes] + [""])
    for idx, tbl in enumerate(s.tables):
        if tbl.rows:
            lines.extend([f"### Table {idx + 1}", "| " + " | ".join(tbl.rows[0]) + " |",
                          "|" + " --- |" * len(tbl.rows[0])] + [f"| {' | '.join(r)} |" for r in tbl.rows[1:]] + [""])
    if s.images:
        lines.extend(["### Images"] + [f"- {i.filename} ({i.width:.1f}x{i.height:.1f} inches)" for i in s.images] + [""])
    if s.hyperlinks:
        lines.extend(["### Hyperlinks"] + [f"- [{h.text}]({h.url})" for h in s.hyperlinks] + [""])
    if s.charts:
        lines.extend(["### Charts"] + [f"- Type: {c.chart_type}" + (f"\n  Title: {c.title}" if c.title else "") for c in s.charts] + [""])
    return "\n".join(lines)


def parse_pptx(pptx_path: str | Path) -> ParsedPresentation:
    pptx_path = Path(pptx_path)
    validate_file(pptx_path)

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        raise PPTParserError(f"无法打开 PPT 文件: {e}")

    result = ParsedPresentation(presentation_title=pptx_path.stem, slide_count=len(prs.slides))
    warnings: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        try:
            sd = SlideData(slide_number=idx)
            layout = slide.slide_layout
            if layout and layout.name:
                sd.layout_name = layout.name
            sd.title, sd.texts, sd.tables, sd.images, sd.hyperlinks, sd.charts = process_slide(slide)
            try:
                sd.notes = extract_notes(slide)
            except Exception as e:
                warnings.append(f"Slide {idx} notes failed: {e}")
            sd.raw_markdown = generate_markdown(sd)
            result.slides.append(sd)
        except Exception as e:
            warnings.append(f"Slide {idx} failed: {e}")
            result.slides.append(SlideData(slide_number=idx, raw_markdown=f"# Slide {idx}\n\n[Parsing Failed: {e}]"))

    result.warnings = warnings
    return result


def to_dict(parsed: ParsedPresentation) -> dict[str, Any]:
    return {
        "presentation_title": parsed.presentation_title,
        "slide_count": parsed.slide_count,
        "slides": [{
            "slide_number": s.slide_number, "title": s.title, "layout_name": s.layout_name,
            "texts": [{"level": t.level, "text": t.text} for t in s.texts],
            "notes": s.notes,
            "tables": [{"rows": t.rows, "row_count": t.row_count, "col_count": t.col_count} for t in s.tables],
            "images": [{"filename": i.filename, "content_type": i.content_type, "width": i.width, "height": i.height} for i in s.images],
            "hyperlinks": [{"text": h.text, "url": h.url} for h in s.hyperlinks],
            "charts": [{"chart_type": c.chart_type, "title": c.title} for c in s.charts],
            "raw_markdown": s.raw_markdown
        } for s in parsed.slides],
        "warnings": parsed.warnings
    }


def to_json(parsed: ParsedPresentation, indent: int = 2) -> str:
    return json.dumps(to_dict(parsed), ensure_ascii=False, indent=indent)


def to_markdown(parsed: ParsedPresentation) -> str:
    lines = [f"# {parsed.presentation_title}", "", f"Total Slides: {parsed.slide_count}", "", "---", ""]
    for slide in parsed.slides:
        lines.extend([slide.raw_markdown, "---", ""])
    if parsed.warnings:
        lines.extend(["## Warnings"] + [f"- {w}" for w in parsed.warnings])
    return "\n".join(lines)


def parse(pptx_path: str | Path, json_output: str | Path | None = None, md_output: str | Path | None = None) -> ParsedPresentation:
    parsed = parse_pptx(pptx_path)
    if json_output:
        Path(json_output).parent.mkdir(parents=True, exist_ok=True)
        Path(json_output).write_text(to_json(parsed), encoding='utf-8')
    if md_output:
        Path(md_output).parent.mkdir(parents=True, exist_ok=True)
        Path(md_output).write_text(to_markdown(parsed), encoding='utf-8')
    return parsed


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("用法: python ppt_parser.py <pptx路径> [输出JSON] [输出Markdown]")
        sys.exit(1)
    try:
        result = parse(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None, sys.argv[3] if len(sys.argv) > 3 else None)
        print(to_json(result))
    except PPTParserError as e:
        print(json.dumps({"status": "error", "reason": str(e)}))
        sys.exit(1)
