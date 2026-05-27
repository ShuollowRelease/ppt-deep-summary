"""
PPT 文本提取模块

使用 python-pptx 从 PPT 文件中提取标题和 bullet list，返回结构化 JSON 数据。
"""

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Paragraph


def extract_bullet_list(text_frame) -> list[dict[str, Any]]:
    """从文本框中提取 bullet list 结构。

    Args:
        text_frame: python-pptx 的 TextFrame 对象

    Returns:
        包含 level 和 text 的字典列表
    """
    bullets: list[dict[str, Any]] = []
    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            bullets.append({
                "level": paragraph.level or 0,
                "text": text
            })
    return bullets


def extract_title_from_shape(shape: Shape) -> str | None:
    """从形状中提取标题文本。

    Args:
        shape: python-pptx 的 Shape 对象

    Returns:
        标题文本，如果不是标题形状则返回 None
    """
    if not shape.has_text_frame:
        return None

    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            return text
    return None


def extract_content_from_shape(shape: Shape) -> dict[str, Any] | None:
    """从形状中提取内容（bullet list）。

    Args:
        shape: python-pptx 的 Shape 对象

    Returns:
        包含形状信息的字典，如果无内容则返回 None
    """
    if not shape.has_text_frame:
        return None

    bullets = extract_bullet_list(shape.text_frame)
    if not bullets:
        return None

    return {
        "shape_type": str(shape.shape_type),
        "bullets": bullets
    }


def process_shape_recursively(shape) -> list[dict[str, Any]]:
    """递归处理形状（包括组形状）。

    Args:
        shape: python-pptx 的 Shape 对象

    Returns:
        提取的内容列表
    """
    results: list[dict[str, Any]] = []

    if isinstance(shape, GroupShape):
        for child_shape in shape.shapes:
            results.extend(process_shape_recursively(child_shape))
    elif shape.has_text_frame:
        content = extract_content_from_shape(shape)
        if content:
            results.append(content)

    return results


def extract_slide_data(slide, slide_index: int) -> dict[str, Any]:
    """从单个幻灯片中提取数据。

    Args:
        slide: python-pptx 的 Slide 对象
        slide_index: 幻灯片索引（从1开始）

    Returns:
        包含幻灯片数据的字典
    """
    slide_data: dict[str, Any] = {
        "slide_number": slide_index,
        "title": None,
        "contents": []
    }

    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title = extract_title_from_shape(shape)
                if title:
                    slide_data["title"] = title
            else:
                contents = process_shape_recursively(shape)
                slide_data["contents"].extend(contents)

    return slide_data


def extract_ppt_text(ppt_path: str | Path) -> dict[str, Any]:
    """从 PPT 文件中提取所有文本内容。

    Args:
        ppt_path: PPT 文件路径

    Returns:
        包含所有幻灯片数据的字典

    Raises:
        FileNotFoundError: 文件不存在
        ValueError: 文件格式无效
        Exception: 其他解析错误
    """
    ppt_path = Path(ppt_path)

    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT 文件不存在: {ppt_path}")

    if ppt_path.suffix.lower() not in ['.pptx', '.ppt']:
        raise ValueError(f"不支持的文件格式: {ppt_path.suffix}")

    try:
        prs = Presentation(str(ppt_path))
    except Exception as e:
        raise Exception(f"无法打开 PPT 文件: {e}")

    result: dict[str, Any] = {
        "file_name": ppt_path.name,
        "total_slides": len(prs.slides),
        "slides": []
    }

    for idx, slide in enumerate(prs.slides, start=1):
        slide_data = extract_slide_data(slide, idx)
        result["slides"].append(slide_data)

    return result


def save_to_json(data: dict[str, Any], output_path: str | Path) -> None:
    """将提取的数据保存为 JSON 文件。

    Args:
        data: 提取的数据字典
        output_path: 输出文件路径
    """
    output_path = Path(output_path)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def main(ppt_path: str, output_path: str | None = None) -> dict[str, Any]:
    """主函数：提取 PPT 文本并返回 JSON。

    Args:
        ppt_path: PPT 文件路径
        output_path: 可选的输出 JSON 文件路径

    Returns:
        提取的数据字典
    """
    result = extract_ppt_text(ppt_path)

    if output_path:
        save_to_json(result, output_path)
        print(f"结果已保存到: {output_path}")

    return result


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法: python extract_text.py <ppt文件路径> [输出JSON路径]")
        sys.exit(1)

    ppt_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    result = main(ppt_file, output_file)
    print(json.dumps(result, ensure_ascii=False, indent=2))
print("hello")